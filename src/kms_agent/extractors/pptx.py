from __future__ import annotations

from hashlib import sha256
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from kms_agent.models import DocumentElement, SourceLocation


class PowerPointExtractor:
    def extract(self, path: str | Path) -> list[DocumentElement]:
        source = Path(path)
        checksum = sha256(source.read_bytes()).hexdigest()
        document_id = checksum[:16]
        deck = Presentation(source)
        elements: list[DocumentElement] = []
        for slide_number, slide in enumerate(deck.slides, start=1):
            for shape in slide.shapes:
                elements.extend(
                    self._extract_shape(
                        shape,
                        slide_number=slide_number,
                        document_id=document_id,
                        title=source.name,
                        checksum=checksum,
                    )
                )
            notes = self._notes_text(slide)
            if notes:
                elements.append(
                    self._element(
                        document_id,
                        source.name,
                        checksum,
                        slide_number,
                        "Speaker Notes",
                        "notes",
                        notes,
                    )
                )
        return elements

    def _extract_shape(
        self,
        shape,
        *,
        slide_number: int,
        document_id: str,
        title: str,
        checksum: str,
    ) -> Iterable[DocumentElement]:
        elements: list[DocumentElement] = []
        name = getattr(shape, "name", "Unnamed Shape")
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                elements.extend(
                    self._extract_shape(
                        child,
                        slide_number=slide_number,
                        document_id=document_id,
                        title=title,
                        checksum=checksum,
                    )
                )
            return elements

        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                elements.append(self._element(document_id, title, checksum, slide_number, name, "text", text))

        if getattr(shape, "has_table", False):
            rows = []
            for row_index, row in enumerate(shape.table.rows, start=1):
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(f"Row {row_index}: " + " | ".join(cells))
            elements.append(
                self._element(
                    document_id,
                    title,
                    checksum,
                    slide_number,
                    f"{name} | Table R1C1:R{len(shape.table.rows)}C{len(shape.table.columns)}",
                    "table",
                    "\n".join(rows),
                )
            )

        if getattr(shape, "has_chart", False):
            chart_text = self._chart_text(shape.chart)
            if chart_text:
                elements.append(self._element(document_id, title, checksum, slide_number, name, "chart", chart_text))

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blob = shape.image.blob
            image_hash = sha256(blob).hexdigest()
            description = (getattr(shape, "alternative_text", "") or "").strip()
            elements.append(
                self._element(
                    document_id,
                    title,
                    checksum,
                    slide_number,
                    name,
                    "image",
                    description or f"Image {name}",
                    metadata={"image_checksum": image_hash, "image_extension": shape.image.ext},
                    binary=blob,
                )
            )
        return elements

    @staticmethod
    def _chart_text(chart) -> str:
        parts: list[str] = []
        try:
            if chart.has_title:
                parts.append("Chart title: " + chart.chart_title.text_frame.text.strip())
        except (AttributeError, ValueError):
            pass
        try:
            for plot_index, plot in enumerate(chart.plots, start=1):
                categories = []
                try:
                    categories = [str(category.label) for category in plot.categories]
                except (AttributeError, TypeError, ValueError):
                    pass
                for series in plot.series:
                    values = list(getattr(series, "values", []) or [])
                    name = getattr(series, "name", f"Series {plot_index}")
                    pairs = [f"{categories[index] if index < len(categories) else index + 1}={value}" for index, value in enumerate(values)]
                    parts.append(f"Series {name}: " + ", ".join(pairs))
        except (AttributeError, TypeError, ValueError):
            pass
        return "\n".join(part for part in parts if part.strip())

    @staticmethod
    def _notes_text(slide) -> str:
        try:
            return slide.notes_slide.notes_text_frame.text.strip()
        except (AttributeError, KeyError, ValueError):
            return ""

    @staticmethod
    def _element(
        document_id: str,
        title: str,
        checksum: str,
        slide_number: int,
        shape_name: str,
        content_type: str,
        text: str,
        metadata: dict | None = None,
        binary: bytes | None = None,
    ) -> DocumentElement:
        location = SourceLocation(
            kind="pptx",
            label=f"Slide {slide_number} | {shape_name}",
            coordinates={"slide_number": slide_number, "shape_name": shape_name},
        )
        return DocumentElement.create(
            document_id=document_id,
            document_title=title,
            document_checksum=checksum,
            location=location,
            content_type=content_type,
            text=text,
            metadata=metadata,
            binary=binary,
        )

