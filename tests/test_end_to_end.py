import tempfile
import unittest
from pathlib import Path

from openpyxl import Workbook
from pptx import Presentation

from tests.helpers import SRC  # noqa: F401
from kms_agent.config import Settings
from kms_agent.service import KMSService


class EndToEndTests(unittest.TestCase):
    def test_powerpoint_and_excel_are_searchable_together(self):
        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            pptx_path = root / "Gateway Design.pptx"
            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[5])
            slide.shapes.title.text = "CAN-FD Gateway"
            slide.notes_slide.notes_text_frame.text = "Diagnostic session timeout is 5000 ms"
            deck.save(pptx_path)

            xlsx_path = root / "Gateway Requirements.xlsx"
            book = Workbook()
            sheet = book.active
            sheet.title = "Timing"
            sheet.append(["ID", "Requirement", "Limit"])
            sheet.append(["GW-TIME-014", "Gateway response", "100 ms"])
            book.save(xlsx_path)

            service = KMSService(settings=Settings(embedding_dimension=64, minimum_score=0.05))
            service.ingest(pptx_path)
            service.ingest(xlsx_path)
            requirement = service.answer("GW-TIME-014 gateway response")
            diagnostic = service.answer("diagnostic session timeout")

        self.assertEqual(requirement.citations[0].document, "Gateway Requirements.xlsx")
        self.assertEqual(diagnostic.citations[0].document, "Gateway Design.pptx")
        self.assertIn("Speaker Notes", diagnostic.citations[0].location)


if __name__ == "__main__":
    unittest.main()

