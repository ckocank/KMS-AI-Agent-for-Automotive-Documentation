import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from tests.helpers import SRC  # noqa: F401
from kms_agent.extractors.pptx import PowerPointExtractor


class PowerPointExtractorTests(unittest.TestCase):
    def test_extracts_slide_shape_and_table_locations(self):
        with tempfile.TemporaryDirectory() as folder:
            path = Path(folder) / "gateway.pptx"
            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[5])
            slide.shapes.title.text = "Gateway Architecture"
            table = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(5), Inches(1)).table
            table.cell(0, 0).text = "Requirement"
            table.cell(0, 1).text = "Limit"
            table.cell(1, 0).text = "GW-TIME-014"
            table.cell(1, 1).text = "100 ms"
            slide.notes_slide.notes_text_frame.text = "Validated on CAN-FD"
            deck.save(path)

            elements = PowerPointExtractor().extract(path)

        labels = [item.location.label for item in elements]
        text = "\n".join(item.text for item in elements)
        self.assertTrue(any("Slide 1" in label and "Title" in label for label in labels))
        self.assertIn("GW-TIME-014", text)
        self.assertIn("Validated on CAN-FD", text)


if __name__ == "__main__":
    unittest.main()

